import io
import os
import tempfile

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse, StreamingResponse
from sqlalchemy.orm import Session

from app.database import get_db
from app.models.creator import Creator, Deck
from app.services.deck_generator import generate_deck, finalize_deck

router = APIRouter(prefix="/api/decks", tags=["decks"])


@router.post("")
def create_deck(
    creator_id: str,
    product_recommendation_id: str,
    actor: str = "internal",
    db: Session = Depends(get_db),
):
    try:
        deck = generate_deck(db, creator_id, product_recommendation_id, actor)
        return _dict(deck)
    except ValueError as e:
        raise HTTPException(400, str(e))


@router.get("/{deck_id}")
def get_deck(deck_id: str, db: Session = Depends(get_db)):
    deck = db.get(Deck, deck_id)
    if not deck:
        raise HTTPException(404, "Deck not found")
    return _dict(deck)


@router.post("/{deck_id}/finalize")
def finalize(deck_id: str, actor: str = "internal", db: Session = Depends(get_db)):
    try:
        deck = finalize_deck(db, deck_id, actor)
        return _dict(deck)
    except ValueError as e:
        raise HTTPException(404, str(e))


@router.get("")
def list_decks(creator_id: str = None, db: Session = Depends(get_db)):
    q = db.query(Deck)
    if creator_id:
        q = q.filter(Deck.creator_id == creator_id)
    return [_dict(d) for d in q.order_by(Deck.created_at.desc()).all()]


@router.get("/{deck_id}/download")
def download_pptx(deck_id: str, db: Session = Depends(get_db)):
    """Generate and stream a .pptx file for the given deck."""
    import httpx as _httpx
    import tempfile
    import os

    deck = db.get(Deck, deck_id)
    if not deck:
        raise HTTPException(404, "Deck not found")

    creator = db.get(Creator, deck.creator_id)
    creator_name = (creator.display_name or creator.handle) if creator else "Creator"
    avatar_url   = (creator.avatar_url or "") if creator else ""
    followers    = creator.follower_count if creator else 0
    niche_tags   = ", ".join(creator.niche or []) if creator else ""

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(500, "python-pptx not installed — run: pip install python-pptx")

    # ── Download creator avatar for embedding ──────────────
    avatar_tmp = None
    if avatar_url:
        try:
            r = _httpx.get(avatar_url, headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                              "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
                "Referer": "https://www.youtube.com/",
            }, timeout=10, follow_redirects=True)
            ct = r.headers.get("content-type", "")
            if r.status_code == 200 and "image" in ct:
                suffix = ".png" if "png" in ct else ".jpg"
                tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
                tmp.write(r.content)
                tmp.close()
                avatar_tmp = tmp.name
        except Exception:
            pass

    # ── Colour tokens ─────────────────────────────────────
    BLACK  = RGBColor(0x05, 0x05, 0x05)
    WHITE  = RGBColor(0xF2, 0xF2, 0xF2)
    GRAY   = RGBColor(0x88, 0x88, 0x88)
    DIM    = RGBColor(0x33, 0x33, 0x33)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slides = deck.slides or []

    # ── Helper: add text box ───────────────────────────────
    def _text(slide, x, y, w, h, text, size, bold=False, color=WHITE, align=None, wrap=True):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.text = text
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Helvetica Neue"
        return box

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = (slide_data.get("type") or "").lower()
        is_cover = (i == 0 or slide_type == "cover")

        # ── Background ────────────────────────────────────
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = BLACK; bg.line.color.rgb = BLACK

        # ── Cover slide: creator photo fills right half ───
        if is_cover and avatar_tmp:
            photo_x = Inches(7.5)
            photo_w = prs.slide_width - Inches(7.5)
            try:
                slide.shapes.add_picture(avatar_tmp, photo_x, Inches(0), photo_w, prs.slide_height)
            except Exception:
                pass
            # gradient overlay strip to blend photo into dark left
            grad = slide.shapes.add_shape(1, Inches(6.8), Inches(0), Inches(1.5), prs.slide_height)
            grad.fill.solid(); grad.fill.fore_color.rgb = BLACK
            grad.fill.fore_color.theme_color  # keep solid for compatibility
            grad.line.color.rgb = BLACK

        # ── Traction/stats slide: avatar top-right badge ──
        elif slide_type in ("traction", "stats") and avatar_tmp:
            badge_size = Inches(1.4)
            try:
                slide.shapes.add_picture(
                    avatar_tmp,
                    prs.slide_width - badge_size - Inches(0.4),
                    Inches(0.15),
                    badge_size, badge_size,
                )
            except Exception:
                pass

        # ── Small avatar watermark on all non-cover slides ─
        elif not is_cover and avatar_tmp:
            mini = Inches(0.48)
            try:
                slide.shapes.add_picture(
                    avatar_tmp,
                    prs.slide_width - mini - Inches(0.35),
                    Inches(0.09),
                    mini, mini,
                )
            except Exception:
                pass

        # ── Top rule ──────────────────────────────────────
        rule = slide.shapes.add_shape(1, Inches(0), Inches(0.6), prs.slide_width, Inches(0.018))
        rule.fill.solid(); rule.fill.fore_color.rgb = DIM; rule.line.color.rgb = DIM

        # ── Slide meta (number + type) ────────────────────
        _text(slide, 0.5, 0.18, 4, 0.35,
              f"{i+1:02d}  {slide_type.upper()}", 8, color=GRAY)

        # ── Text area width: narrow on cover (photo takes right half) ─
        text_w = 6.8 if is_cover else 12.0

        # ── Headline ──────────────────────────────────────
        headline = slide_data.get("headline") or slide_data.get("title") or ""
        headline_size = 42 if is_cover else 34
        _text(slide, 0.5, 1.1, text_w, 1.6, headline, headline_size, bold=True)

        # ── Cover extras: creator name + niche tag below headline ──
        if is_cover:
            sub = creator_name
            if niche_tags:
                sub += f"  ·  {niche_tags}"
            _text(slide, 0.5, 2.85, text_w, 0.4, sub, 13, color=GRAY)
            if followers:
                def _fmt_num(n):
                    if n >= 1_000_000: return f"{n/1_000_000:.1f}M".replace(".0M","M")
                    if n >= 1_000: return f"{n//1_000}K"
                    return str(n)
                _text(slide, 0.5, 3.3, text_w, 0.4,
                      f"{_fmt_num(followers)} subscribers", 11, color=GRAY)

        # ── Body / Bullets ────────────────────────────────
        body_text = slide_data.get("body") or ""
        bullets   = slide_data.get("bullets") or []

        if body_text:
            _text(slide, 0.5, 2.9, text_w, 3.5, body_text, 15, color=GRAY)

        if bullets:
            bul_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.1), Inches(text_w), Inches(3.3)
            )
            tf4 = bul_box.text_frame; tf4.word_wrap = True
            for j, bullet in enumerate(bullets[:6]):
                p4 = tf4.paragraphs[0] if j == 0 else tf4.add_paragraph()
                p4.text = f"→  {bullet}"
                run4 = p4.runs[0]
                run4.font.size = Pt(14); run4.font.color.rgb = GRAY
                run4.font.name = "Helvetica Neue"
                p4.space_before = Pt(5)

        # ── Watermark ─────────────────────────────────────
        _text(slide, 9.0, 7.0, 4.2, 0.35,
              f"Creator Forge  ·  {creator_name}", 7.5,
              color=DIM, align=PP_ALIGN.RIGHT)

    # ── Stream ────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    if avatar_tmp:
        try:
            os.unlink(avatar_tmp)
        except Exception:
            pass

    filename = f"{creator_name.replace(' ','_')}_pitch_deck.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def _dict(d: Deck) -> dict:
    return {
        "id": d.id, "creator_id": d.creator_id,
        "product_recommendation_id": d.product_recommendation_id,
        "title": d.title, "slides": d.slides, "version": d.version,
        "status": d.status,
        "created_at": d.created_at.isoformat() if d.created_at else None,
        "updated_at": d.updated_at.isoformat() if d.updated_at else None,
    }
